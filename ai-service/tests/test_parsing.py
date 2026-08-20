import pytest

from app.core.parsing import parse_file


def test_parse_txt_utf8(tmp_path):
    f = tmp_path / "a.txt"
    f.write_text("你好世界", encoding="utf-8")
    pages = parse_file(str(f))
    assert pages[0][0] == "你好世界"
    assert pages[0][1] == 1


def test_parse_txt_gbk_fallback(tmp_path):
    f = tmp_path / "b.txt"
    f.write_text("中文内容", encoding="gbk")
    pages = parse_file(str(f))
    assert pages[0][0] == "中文内容"


def test_parse_legacy_doc_rejected(tmp_path):
    f = tmp_path / "c.doc"
    f.write_text("x")
    with pytest.raises(ValueError, match="旧版"):
        parse_file(str(f))


def test_parse_unsupported_extension_rejected(tmp_path):
    f = tmp_path / "d.exe"
    f.write_text("x")
    with pytest.raises(ValueError, match="不支持"):
        parse_file(str(f))


def test_parse_markdown(tmp_path):
    f = tmp_path / "e.md"
    f.write_text("# 标题\n\n这是内容", encoding="utf-8")
    pages = parse_file(str(f))
    assert "标题" in pages[0][0]
    assert "这是内容" in pages[0][0]


def test_parse_xlsx(tmp_path):
    import openpyxl

    f = tmp_path / "f.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["姓名", "部门"])
    ws.append(["张三", "研发部"])
    wb.save(f)

    pages = parse_file(str(f))
    assert len(pages) >= 1
    assert "张三" in pages[0][0]
    assert "研发部" in pages[0][0]


def test_parse_pptx(tmp_path):
    from pptx import Presentation

    f = tmp_path / "g.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "会议纪要"
    prs.save(f)

    pages = parse_file(str(f))
    assert any("会议纪要" in p[0] for p in pages)
