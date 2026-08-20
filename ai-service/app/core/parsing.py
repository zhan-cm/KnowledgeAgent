from pathlib import Path


def parse_file(file_path):
    """解析文档，返回 [(页文本, 页码)]，页码从 1 开始"""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    if ext == ".docx":
        return _parse_docx(file_path)
    if ext == ".doc":
        raise ValueError("暂不支持旧版 .doc 格式，请转换为 .docx 后重新上传")
    if ext in (".txt", ".md", ".markdown"):
        return _parse_txt(file_path)
    if ext == ".xlsx":
        return _parse_xlsx(file_path)
    if ext == ".pptx":
        return _parse_pptx(file_path)
    raise ValueError(f"不支持的文件类型: {ext}")


def _parse_pdf(path):
    import fitz

    pages = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                pages.append((text, i + 1))
    return pages


def _parse_docx(path):
    import docx

    document = docx.Document(path)
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [(text, 1)] if text.strip() else []


def _parse_txt(path):
    for encoding in ("utf-8", "gbk", "utf-16"):
        try:
            text = Path(path).read_text(encoding=encoding)
            return [(text, 1)] if text.strip() else []
        except UnicodeDecodeError:
            continue
    raise ValueError("无法识别文本文件编码")


def _parse_xlsx(path):
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    try:
        for i, sheet in enumerate(workbook.worksheets, 1):
            lines = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append(" | ".join(cells))
            if lines:
                pages.append(("\n".join(lines), i))
    finally:
        workbook.close()
    return pages


def _parse_pptx(path):
    from pptx import Presentation

    presentation = Presentation(path)
    pages = []
    for i, slide in enumerate(presentation.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        if texts:
            pages.append(("\n".join(texts), i))
    return pages
